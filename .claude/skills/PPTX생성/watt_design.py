"""
WATT Design System — IBM Carbon 기반 PPTX 디자인 라이브러리

와트 브랜드 정체성을 모든 PPTX에 일관 적용.
참조: .claude/design_md/ibm/DESIGN.md + 와트 홈페이지 (www.wattsolution.co.kr)

사용 예시:
    from watt_design import Deck
    from components import slide_cover, slide_contents, slide_card_3col

    deck = Deck(total_pages=12)
    slide_cover(deck, title="...", subtitle="...", tag="...")
    ...
    deck.save("제안서.pptx")
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ═════════════════════════════════════════════════════════
# 디자인 토큰 (IBM Carbon)
# ═════════════════════════════════════════════════════════
NAVY = RGBColor(0x16, 0x16, 0x16)       # Gray 100 — 주 텍스트·다크 서피스
CYAN = RGBColor(0x0F, 0x62, 0xFE)       # IBM Blue 60 — 유일 액센트
DARK = RGBColor(0x16, 0x16, 0x16)       # Gray 100 — body text
GRAY = RGBColor(0x52, 0x52, 0x52)       # Gray 70 — secondary text
LIGHT_GRAY = RGBColor(0xE0, 0xE0, 0xE0) # Gray 20 — borders
SECTION_BG = RGBColor(0xF4, 0xF4, 0xF4) # Gray 10 — layer-01
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_RED = RGBColor(0xDA, 0x1E, 0x28)   # Red 60
ACCENT_GREEN = RGBColor(0x24, 0xA1, 0x48) # Green 50
BLUE_10 = RGBColor(0xED, 0xF5, 0xFF)
BLUE_70 = RGBColor(0x00, 0x43, 0xCE)
FONT = "맑은 고딕"


class Deck:
    """16:9 프레젠테이션 래퍼. 페이지 번호는 자동 카운트."""

    def __init__(self, total_pages: int, brand_name: str = "㈜와트  |  WATT CO. LTD."):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.SW = self.prs.slide_width
        self.SH = self.prs.slide_height
        self.blank = self.prs.slide_layouts[6]
        self.total_pages = total_pages
        self.brand_name = brand_name
        self._page_counter = 0

    def add_slide(self):
        self._page_counter += 1
        s = self.prs.slides.add_slide(self.blank)
        s._watt_page = self._page_counter  # 커스텀 속성
        return s

    @property
    def current_page(self):
        return self._page_counter

    def save(self, path):
        self.prs.save(path)


# ═════════════════════════════════════════════════════════
# Primitive Helpers
# ═════════════════════════════════════════════════════════
def add_rect(slide, x, y, w, h, fill=None, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.shadow.inherit = False
    if fill is not None:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    return shp


def add_text(slide, x, y, w, h, text, *,
             size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, font=FONT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=14, color=DARK, bullet="•  "):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
        r = p.add_run()
        if item == "":
            r.text = ""
        elif item.lstrip().startswith("•"):
            r.text = item
        else:
            r.text = bullet + item
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb


# ═════════════════════════════════════════════════════════
# Page Frame (페이지 공통 요소 — 사이드바·타이틀·섹션태그·페이지번호·브랜드)
# ═════════════════════════════════════════════════════════
def page_frame(deck, slide, title, subtitle=None, section_tag=None):
    add_rect(slide, 0, 0, deck.SW, deck.SH, fill=WHITE)
    add_rect(slide, 0, 0, Inches(0.35), deck.SH, fill=NAVY)
    if section_tag:
        add_text(slide, Inches(0.7), Inches(0.35), Inches(11), Inches(0.3),
                 section_tag, size=11, bold=True, color=CYAN)
        add_text(slide, Inches(0.7), Inches(0.65), Inches(11), Inches(0.55),
                 title, size=22, bold=True, color=NAVY)
        add_rect(slide, Inches(0.7), Inches(1.2), Inches(0.6), Emu(38100), fill=CYAN)
    else:
        add_text(slide, Inches(0.7), Inches(0.45), Inches(11), Inches(0.55),
                 title, size=24, bold=True, color=NAVY)
        add_rect(slide, Inches(0.7), Inches(1.05), Inches(0.6), Emu(38100), fill=CYAN)
    if subtitle:
        add_text(slide, Inches(0.7), Inches(1.3), Inches(11), Inches(0.4),
                 subtitle, size=13, color=GRAY)
    # 페이지 번호 (우하단)
    page = getattr(slide, "_watt_page", deck.current_page)
    add_text(slide, Inches(12.3), Inches(7.05), Inches(0.9), Inches(0.3),
             f"{page} / {deck.total_pages}", size=10, color=GRAY, align=PP_ALIGN.RIGHT)
    # 브랜드 (좌하단)
    add_text(slide, Inches(0.7), Inches(7.05), Inches(6), Inches(0.3),
             deck.brand_name, size=10, color=GRAY)


# ═════════════════════════════════════════════════════════
# Card Primitive (섹션 BG + 상단 accent bar + 타이틀 + 불릿)
# ═════════════════════════════════════════════════════════
def card(slide, x, y, w, h, title, body_items, *, accent=CYAN):
    add_rect(slide, x, y, w, h, fill=SECTION_BG, line=LIGHT_GRAY)
    add_rect(slide, x, y, w, Emu(76200), fill=accent)
    add_text(slide, x + Inches(0.25), y + Inches(0.2), w - Inches(0.5), Inches(0.45),
             title, size=15, bold=True, color=NAVY)
    add_bullets(slide, x + Inches(0.25), y + Inches(0.75),
                w - Inches(0.5), h - Inches(0.9), body_items, size=12, color=DARK)


def kpi_box(slide, x, y, w, h, big_number, unit, label, *, accent=CYAN):
    add_rect(slide, x, y, w, h, fill=WHITE, line=LIGHT_GRAY)
    add_rect(slide, x, y + h - Emu(76200), w, Emu(76200), fill=accent)
    add_text(slide, x, y + Inches(0.25), w, Inches(1.3),
             big_number, size=56, bold=True, color=accent,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, x, y + Inches(1.55), w, Inches(0.35),
             unit, size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(slide, x, y + Inches(1.95), w, Inches(0.9),
             label, size=12, color=GRAY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)


__all__ = [
    "Deck", "page_frame", "card", "kpi_box",
    "add_rect", "add_text", "add_bullets",
    "NAVY", "CYAN", "DARK", "GRAY", "LIGHT_GRAY", "SECTION_BG", "WHITE",
    "ACCENT_RED", "ACCENT_GREEN", "BLUE_10", "BLUE_70", "FONT",
    "Inches", "Pt", "Emu", "MSO_SHAPE", "PP_ALIGN", "MSO_ANCHOR",
]
